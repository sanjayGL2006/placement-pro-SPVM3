# generate_ppt.py
# -------------------------------------------------
# Creates a simple PowerPoint presentation for the
# placement‑pro project using python-pptx.
# -------------------------------------------------
from pathlib import Path
# pyrefly: ignore [missing-import]
from pptx import Presentation
# pyrefly: ignore [missing-import]
from pptx.util import Inches, Pt

# -----------------------------------------------------------------
# Configuration – update these sections with your own content
# -----------------------------------------------------------------
PROJECT_NAME = "Placement‑Pro"
SLIDES = [
    {
        "title": "Placement‑Pro Overview",
        "bullets": [
            "Web‑based placement portal for students",
            "Flask backend with PostgreSQL",
            "RESTful APIs for notifications, recycle‑bin, AI integration",
            "Docker‑ready for easy deployment",
        ],
    },
    {
        "title": "Key Modules",
        "bullets": [
            "🔹 backend/database.py – DB connection pool",
            "🔹 backend/routes/ – API endpoints",
            "🔹 backend/scheduler.py – background tasks",
            "🔹 frontend/ – (future) UI layer",
        ],
    },
    {
        "title": "Architecture Diagram",
        "bullets": [],
        "image_path": "architecture.png",  # optional image in same folder
    },
    {
        "title": "Next Steps",
        "bullets": [
            "Add unit tests",
            "Integrate Firebase for auth & hosting",
            "Create CI/CD pipeline",
        ],
    },
]

def add_slide(prs, title, bullets=None, image_path=None):
    slide_layout = prs.slide_layouts[5]  # Title & Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if bullets:
        tf = slide.placeholders[1].text_frame
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
    if image_path:
        img_path = Path(image_path)
        if img_path.is_file():
            left = Inches(1)
            top = Inches(2)
            slide.shapes.add_picture(str(img_path), left, top, width=Inches(6))
        else:
            print(f"⚠️ Image not found: {img_path}")

def main():
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = f"{PROJECT_NAME} Presentation"
    slide.placeholders[1].text = "Generated with python‑pptx"
    # Add slides
    for s in SLIDES:
        add_slide(prs, s["title"], s.get("bullets"), s.get("image_path"))
    out_path = Path(__file__).parent / f"{PROJECT_NAME}.pptx"
    prs.save(out_path)
    print(f"✅ Presentation saved to: {out_path}")

if __name__ == "__main__":
    main()
